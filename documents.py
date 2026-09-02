import sys
import json
import base64
import io
import requests

from reportlab.lib.pagesizes import A4, LETTER, LEGAL
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, BaseDocTemplate, PageTemplate, Frame,
    Paragraph, Spacer, Table, TableStyle, HRFlowable, Image,
    PageBreak, FrameBreak, KeepTogether
)
from reportlab.platypus.tableofcontents import TableOfContents
from reportlab.platypus.flowables import Flowable
from reportlab.graphics.shapes import (
    Drawing, Rect, Circle, Ellipse, Line, Polygon, PolyLine, String, Wedge
)
from reportlab.graphics import renderPM
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.legends import Legend
from reportlab.graphics.barcode.qr import QrCodeWidget
from reportlab.graphics.barcode import code128

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference, Series
from openpyxl.formatting.rule import ColorScaleRule, DataBarRule
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table as XlsxTable, TableStyleInfo
from openpyxl.drawing.image import Image as XlsxImage

from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

COR_PRIMARIA = colors.HexColor("#0F172A")
COR_SECUNDARIA = colors.HexColor("#4F46E5")
COR_DESTAQUE = colors.HexColor("#F59E0B")
COR_TEXTO = colors.HexColor("#475569")
COR_CINZA_CLARO = colors.HexColor("#F8FAFC")
PALETA = ["#4F46E5", "#0EA5E9", "#F59E0B", "#EF4444", "#10B981", "#8B5CF6", "#EC4899", "#14B8A6"]
PALETA_SEM_HASH = [c.lstrip("#") for c in PALETA]

TAMANHOS_PAGINA = {"A4": A4, "LETTER": LETTER, "LEGAL": LEGAL}
ALINHAMENTOS = {"left": "LEFT", "center": "CENTER", "right": "RIGHT", "decimal": "DECIMAL"}


# ══════════════════════════════════════════════════════════════════
# PDF — estilos e helpers
# ══════════════════════════════════════════════════════════════════

def get_styles():
    styles = getSampleStyleSheet()
    return {
        "titulo": ParagraphStyle("Titulo", parent=styles["Title"], fontName="Helvetica-Bold",
                                  fontSize=26, leading=30, textColor=COR_PRIMARIA, spaceAfter=10),
        "subtitulo": ParagraphStyle("Subtitulo", parent=styles["Normal"], fontName="Helvetica",
                                     fontSize=13, leading=17, textColor=COR_SECUNDARIA, spaceAfter=18),
        "heading": ParagraphStyle("Heading", parent=styles["Heading1"], fontName="Helvetica-Bold",
                                   fontSize=15, leading=19, textColor=COR_PRIMARIA,
                                   spaceBefore=16, spaceAfter=8),
        "heading2": ParagraphStyle("Heading2", parent=styles["Heading2"], fontName="Helvetica-Bold",
                                    fontSize=12.5, leading=16, textColor=COR_SECUNDARIA,
                                    spaceBefore=10, spaceAfter=6),
        "heading3": ParagraphStyle("Heading3", parent=styles["Heading2"], fontName="Helvetica-Bold",
                                    fontSize=11, leading=14, textColor=COR_TEXTO,
                                    spaceBefore=8, spaceAfter=4),
        "corpo": ParagraphStyle("Corpo", parent=styles["Normal"], fontName="Helvetica",
                                 fontSize=10.5, leading=15.5, textColor=COR_TEXTO, spaceAfter=8),
        "bullet": ParagraphStyle("Bullet", parent=styles["Normal"], fontName="Helvetica",
                                  fontSize=10.5, leading=15, leftIndent=14, spaceAfter=5,
                                  textColor=COR_TEXTO),
        "toc_titulo": ParagraphStyle("TocTitulo", parent=styles["Heading1"], fontName="Helvetica-Bold",
                                      fontSize=18, textColor=COR_PRIMARIA, spaceAfter=14),
        "callout": ParagraphStyle("Callout", parent=styles["Normal"], fontName="Helvetica",
                                   fontSize=10, leading=14, textColor=COR_PRIMARIA),
        "celula": ParagraphStyle("Celula", parent=styles["Normal"], fontName="Helvetica",
                                  fontSize=9.5, leading=13, textColor=COR_TEXTO),
        "citacao": ParagraphStyle("Citacao", parent=styles["Normal"], fontName="Helvetica-Oblique",
                                   fontSize=11.5, leading=17, textColor=COR_SECUNDARIA,
                                   leftIndent=22, spaceAfter=10, spaceBefore=6),
        "capa_rodape": ParagraphStyle("CapaRodape", parent=styles["Normal"], fontName="Helvetica",
                                       fontSize=9.5, leading=13, textColor=colors.white, alignment=1),
        "capa_titulo": ParagraphStyle("CapaTitulo", parent=styles["Title"], fontName="Helvetica-Bold",
                                       fontSize=30, leading=36, textColor=colors.white, alignment=1),
        "capa_subtitulo": ParagraphStyle("CapaSubtitulo", parent=styles["Normal"], fontName="Helvetica",
                                          fontSize=14, leading=19, textColor=colors.white, alignment=1,
                                          spaceBefore=10),
    }


def get_toc_styles():
    return [
        ParagraphStyle("TOCHeading1", fontName="Helvetica-Bold", fontSize=11, leading=16,
                        textColor=COR_PRIMARIA, leftIndent=0),
        ParagraphStyle("TOCHeading2", fontName="Helvetica", fontSize=10, leading=14,
                        textColor=COR_TEXTO, leftIndent=14),
        ParagraphStyle("TOCHeading3", fontName="Helvetica", fontSize=9, leading=12,
                        textColor=colors.grey, leftIndent=26),
    ]


def linha_divisoria(cor_hex=None, espessura=1.5):
    cor = colors.HexColor(cor_hex) if cor_hex else COR_DESTAQUE
    return HRFlowable(width="100%", thickness=espessura, color=cor, spaceBefore=4, spaceAfter=14)


def rodape_pagina(canvas_obj, doc, texto_custom=None, cor_hex=None):
    canvas_obj.saveState()
    canvas_obj.setStrokeColor(COR_CINZA_CLARO)
    canvas_obj.line(2 * cm, 1.7 * cm, doc.pagesize[0] - 2 * cm, 1.7 * cm)
    canvas_obj.setFont("Helvetica", 8)
    canvas_obj.setFillColor(colors.HexColor(cor_hex) if cor_hex else colors.grey)
    rodape_texto = texto_custom if texto_custom else f"Página {doc.page}"
    rodape_texto = rodape_texto.replace("{page}", str(doc.page))
    canvas_obj.drawRightString(doc.pagesize[0] - 2 * cm, 1.2 * cm, rodape_texto)
    canvas_obj.restoreState()


def desenhar_marca_dagua(canvas_obj, doc, watermark):
    canvas_obj.saveState()
    canvas_obj.setFont("Helvetica-Bold", watermark.get("font_size", 60))
    cor = colors.HexColor(watermark.get("color", "#EEEEEE"))
    canvas_obj.setFillColor(cor)
    canvas_obj.setFillAlpha(watermark.get("opacity", 0.3))
    canvas_obj.translate(doc.pagesize[0] / 2, doc.pagesize[1] / 2)
    canvas_obj.rotate(watermark.get("angle", 45))
    canvas_obj.drawCentredString(0, 0, watermark.get("text", ""))
    canvas_obj.restoreState()


def desenhar_faixa_lateral(canvas_obj, doc, cor_hex, largura_pt=14):
    canvas_obj.saveState()
    canvas_obj.setFillColor(colors.HexColor(cor_hex))
    canvas_obj.rect(0, 0, largura_pt, doc.pagesize[1], fill=1, stroke=0)
    canvas_obj.restoreState()


def rodape_e_extras_pagina(header_text, watermark, footer_spec=None, faixa_lateral_hex=None, header_cor_hex=None, capa_spec=None):
    def _draw(canvas_obj, doc):
        if capa_spec and doc.page == 1:
            desenhar_fundo_capa(canvas_obj, doc, capa_spec)
            return
        if watermark:
            desenhar_marca_dagua(canvas_obj, doc, watermark)
        if faixa_lateral_hex:
            desenhar_faixa_lateral(canvas_obj, doc, faixa_lateral_hex)
        texto_rodape = footer_spec.get("text") if footer_spec else None
        cor_rodape = footer_spec.get("color") if footer_spec else None
        rodape_pagina(canvas_obj, doc, texto_rodape, cor_rodape)
        if header_text:
            canvas_obj.saveState()
            canvas_obj.setFont("Helvetica-Bold", 8)
            canvas_obj.setFillColor(colors.HexColor(header_cor_hex) if header_cor_hex else COR_SECUNDARIA)
            canvas_obj.drawString(2 * cm, doc.pagesize[1] - 1.4 * cm, header_text)
            canvas_obj.setStrokeColor(COR_CINZA_CLARO)
            canvas_obj.line(2 * cm, doc.pagesize[1] - 1.6 * cm, doc.pagesize[0] - 2 * cm, doc.pagesize[1] - 1.6 * cm)
            canvas_obj.restoreState()
    return _draw


def baixar_imagem(url, max_width_cm=15, max_height_cm=None):
    try:
        resp = requests.get(url, timeout=8)
        resp.raise_for_status()
        img = Image(io.BytesIO(resp.content))
        largura_max = max_width_cm * cm
        proporcao = img.imageHeight / img.imageWidth
        img.drawWidth = largura_max
        img.drawHeight = largura_max * proporcao
        if max_height_cm and img.drawHeight > max_height_cm * cm:
            img.drawHeight = max_height_cm * cm
            img.drawWidth = img.drawHeight / proporcao
        return img
    except Exception:
        return None


def baixar_bytes_imagem(url):
    try:
        resp = requests.get(url, timeout=8)
        resp.raise_for_status()
        return io.BytesIO(resp.content)
    except Exception:
        return None


def montar_tabela(t):
    """Constrói uma Table completa aceitando tudo: merge de células
    (span), alinhamento por coluna, cor de fundo/texto por célula
    individual, imagem dentro de célula, bordas seletivas, larguras
    de coluna, cabeçalho customizável, linhas zebra customizáveis.
    Campos além de headers/rows são todos opcionais."""
    headers = t["headers"]
    rows = t["rows"]
    dados = [headers] + [list(r) for r in rows]

    for i, linha in enumerate(dados):
        for j, valor in enumerate(linha):
            if isinstance(valor, dict) and valor.get("image_url"):
                img = baixar_imagem(valor["image_url"], max_width_cm=3)
                dados[i][j] = img if img else ""
            elif isinstance(valor, dict):
                dados[i][j] = valor.get("text", "")

    largura_disponivel = A4[0] - 4 * cm
    if t.get("col_widths_cm"):
        larguras = [w * cm for w in t["col_widths_cm"]]
    else:
        largura_col = largura_disponivel / len(headers)
        larguras = [largura_col] * len(headers)

    tabela = Table(dados, colWidths=larguras, repeatRows=1)

    cor_cabecalho_bg = colors.HexColor(t.get("header_bg", "#0F172A"))
    cor_cabecalho_texto = colors.HexColor(t.get("header_text_color", "#FFFFFF"))
    cor_zebra_1 = colors.HexColor(t.get("zebra_color_1", "#FFFFFF"))
    cor_zebra_2 = colors.HexColor(t.get("zebra_color_2", "#F8FAFC"))

    estilo = [
        ("BACKGROUND", (0, 0), (-1, 0), cor_cabecalho_bg),
        ("TEXTCOLOR", (0, 0), (-1, 0), cor_cabecalho_texto),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), t.get("font_size", 9.5)),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [cor_zebra_1, cor_zebra_2]),
        ("TOPPADDING", (0, 0), (-1, -1), t.get("cell_padding", 7)),
        ("BOTTOMPADDING", (0, 0), (-1, -1), t.get("cell_padding", 7)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]

    if t.get("borders"):
        for b in t["borders"]:
            tipo, inicio, fim, esp = b[0], tuple(b[1]), tuple(b[2]), b[3]
            cor_borda = colors.HexColor(b[4]) if len(b) > 4 else colors.HexColor("#DDDDDD")
            estilo.append((tipo, inicio, fim, esp, cor_borda))
    else:
        estilo.append(("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#DDDDDD")))

    if t.get("align"):
        for col_idx, alinhamento in t["align"].items():
            col = int(col_idx)
            estilo.append(("ALIGN", (col, 0), (col, -1), ALINHAMENTOS.get(alinhamento, "LEFT")))

    if t.get("merge"):
        for span in t["merge"]:
            estilo.append(("SPAN", tuple(span[0]), tuple(span[1])))

    if t.get("cell_styles"):
        for cs in t["cell_styles"]:
            pos = (cs["col"], cs["row"])
            if cs.get("bg"):
                estilo.append(("BACKGROUND", pos, pos, colors.HexColor(cs["bg"])))
            if cs.get("text_color"):
                estilo.append(("TEXTCOLOR", pos, pos, colors.HexColor(cs["text_color"])))
            if cs.get("bold"):
                estilo.append(("FONTNAME", pos, pos, "Helvetica-Bold"))
            if cs.get("align"):
                estilo.append(("ALIGN", pos, pos, ALINHAMENTOS.get(cs["align"], "LEFT")))

    tabela.setStyle(TableStyle(estilo))
    return tabela


def montar_runs_paragrafo(runs, style):
    partes = []
    for r in runs:
        texto = r.get("text", "")
        if r.get("bold"):
            texto = f"<b>{texto}</b>"
        if r.get("italic"):
            texto = f"<i>{texto}</i>"
        if r.get("underline"):
            texto = f"<u>{texto}</u>"
        if r.get("strike"):
            texto = f"<strike>{texto}</strike>"
        if r.get("superscript"):
            texto = f"<super>{texto}</super>"
        if r.get("subscript"):
            texto = f"<sub>{texto}</sub>"
        if r.get("color"):
            texto = f'<font color="{r["color"]}">{texto}</font>'
        if r.get("font_size"):
            texto = f'<font size="{r["font_size"]}">{texto}</font>'
        if r.get("link"):
            texto = f'<link href="{r["link"]}" color="{COR_SECUNDARIA.hexval()}">{texto}</link>'
        partes.append(texto)
    return Paragraph("".join(partes), style)


def montar_grafico(chart_spec):
    tipo = chart_spec.get("type", "bar")
    labels = chart_spec.get("labels", [])
    series = chart_spec.get("series", [])
    largura = chart_spec.get("width", 420)
    altura = chart_spec.get("height", 220)

    drawing = Drawing(largura, altura)

    if chart_spec.get("title"):
        drawing.add(String(largura / 2, altura - 12, chart_spec["title"],
                            fontName="Helvetica-Bold", fontSize=11,
                            fillColor=COR_PRIMARIA, textAnchor="middle"))

    area_top = altura - 30
    cores = [colors.HexColor(c) for c in (chart_spec.get("colors") or PALETA)]

    if tipo == "pie":
        pie = Pie()
        pie.x = largura / 2 - 70
        pie.y = 10
        pie.width = 140
        pie.height = 140
        pie.data = series[0]["data"] if series else []
        pie.labels = labels
        pie.slices.strokeWidth = 0.5
        pie.slices.strokeColor = colors.white
        if chart_spec.get("donut"):
            pie.innerRadiusFraction = 0.55
        for i in range(len(pie.data)):
            pie.slices[i].fillColor = cores[i % len(cores)]
        drawing.add(pie)

    elif tipo == "line":
        chart = HorizontalLineChart()
        chart.x = 50
        chart.y = 30
        chart.width = largura - 90
        chart.height = area_top - 40
        chart.data = [s["data"] for s in series]
        chart.categoryAxis.categoryNames = labels
        chart.categoryAxis.labels.fontSize = 7
        chart.valueAxis.labels.fontSize = 7
        for i, s in enumerate(series):
            chart.lines[i].strokeColor = cores[i % len(cores)]
            chart.lines[i].strokeWidth = 2.2
            if chart_spec.get("show_markers"):
                chart.lines[i].symbol = None
        drawing.add(chart)
        if len(series) > 1:
            legend = Legend()
            legend.x = largura - 5
            legend.y = area_top - 10
            legend.colorNamePairs = [(cores[i % len(cores)], s.get("name", f"Série {i+1}")) for i, s in enumerate(series)]
            legend.fontSize = 7
            legend.alignment = "right"
            drawing.add(legend)

    else:  # bar (default)
        chart = VerticalBarChart()
        chart.x = 50
        chart.y = 30
        chart.width = largura - 90
        chart.height = area_top - 40
        chart.data = [s["data"] for s in series]
        chart.categoryAxis.categoryNames = labels
        chart.categoryAxis.labels.fontSize = 7
        chart.valueAxis.labels.fontSize = 7
        chart.groupSpacing = chart_spec.get("group_spacing", 10)
        chart.barSpacing = chart_spec.get("bar_spacing", 2)
        if chart_spec.get("horizontal"):
            chart.categoryAxis.labels.angle = 0
        for i in range(len(series)):
            chart.bars[i].fillColor = cores[i % len(cores)]
        drawing.add(chart)
        if len(series) > 1:
            legend = Legend()
            legend.x = largura - 5
            legend.y = area_top - 10
            legend.colorNamePairs = [(cores[i % len(cores)], s.get("name", f"Série {i+1}")) for i, s in enumerate(series)]
            legend.fontSize = 7
            legend.alignment = "right"
            drawing.add(legend)

    return drawing


def montar_formas_livres(shapes_spec):
    largura = shapes_spec.get("width", 420)
    altura = shapes_spec.get("height", 200)
    drawing = Drawing(largura, altura)

    for f in shapes_spec.get("shapes", []):
        tipo = f.get("type")
        cor_preenchimento = colors.HexColor(f["fill"]) if f.get("fill") else None
        cor_borda = colors.HexColor(f["stroke"]) if f.get("stroke") else None
        largura_borda = f.get("stroke_width", 1)

        if tipo == "rect":
            drawing.add(Rect(f["x"], f["y"], f["w"], f["h"],
                              fillColor=cor_preenchimento, strokeColor=cor_borda,
                              strokeWidth=largura_borda, rx=f.get("radius", 0), ry=f.get("radius", 0)))
        elif tipo == "circle":
            drawing.add(Circle(f["cx"], f["cy"], f["r"],
                                fillColor=cor_preenchimento, strokeColor=cor_borda, strokeWidth=largura_borda))
        elif tipo == "ellipse":
            drawing.add(Ellipse(f["cx"], f["cy"], f["rx"], f["ry"],
                                 fillColor=cor_preenchimento, strokeColor=cor_borda, strokeWidth=largura_borda))
        elif tipo == "line":
            drawing.add(Line(f["x1"], f["y1"], f["x2"], f["y2"],
                              strokeColor=cor_borda or colors.black, strokeWidth=largura_borda))
        elif tipo == "polygon":
            pontos = []
            for pt in f["points"]:
                pontos.extend(pt)
            drawing.add(Polygon(pontos, fillColor=cor_preenchimento, strokeColor=cor_borda, strokeWidth=largura_borda))
        elif tipo == "polyline":
            pontos = []
            for pt in f["points"]:
                pontos.extend(pt)
            drawing.add(PolyLine(pontos, strokeColor=cor_borda or colors.black, strokeWidth=largura_borda))
        elif tipo == "text":
            drawing.add(String(f["x"], f["y"], f.get("text", ""),
                                fontName=f.get("font", "Helvetica-Bold"),
                                fontSize=f.get("font_size", 12),
                                fillColor=cor_preenchimento or colors.black,
                                textAnchor=f.get("align", "start")))
        elif tipo == "wedge":
            drawing.add(Wedge(f["cx"], f["cy"], f["r"], f["start_angle"], f["end_angle"],
                               fillColor=cor_preenchimento, strokeColor=cor_borda, strokeWidth=largura_borda))

    return drawing


def montar_qrcode(spec):
    largura = spec.get("size", 100)
    widget = QrCodeWidget(spec.get("content", ""))
    b = widget.getBounds()
    w_original = b[2] - b[0]
    h_original = b[3] - b[1]
    drawing = Drawing(largura, largura, transform=[largura / w_original, 0, 0, largura / h_original, 0, 0])
    drawing.add(widget.draw())
    return drawing


def montar_barcode128(spec):
    return code128.Code128(spec.get("content", ""), barHeight=spec.get("height", 15) * 1.0, barWidth=spec.get("bar_width", 0.7))


def montar_callout(callout_spec, style):
    texto = callout_spec.get("text", "")
    icone = callout_spec.get("icon", "")
    if icone:
        texto = f"{icone}&nbsp;&nbsp;{texto}"
    cor_fundo = colors.HexColor(callout_spec.get("bg", "#EEF2FF"))
    cor_borda = colors.HexColor(callout_spec.get("border_color", "#4F46E5"))
    p = Paragraph(texto, style)
    tabela = Table([[p]], colWidths=[A4[0] - 4 * cm])
    tabela.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), cor_fundo),
        ("LINEBEFORE", (0, 0), (0, -1), 3, cor_borda),
        ("TOPPADDING", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING", (0, 0), (-1, -1), 14),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
    ]))
    return tabela


def montar_citacao(quote_spec, style):
    texto = quote_spec.get("text", "")
    autor = quote_spec.get("author", "")
    itens = [Paragraph(f'"{texto}"', style)]
    if autor:
        itens.append(Paragraph(f"— {autor}", ParagraphStyle(
            "AutorCitacao", parent=style, fontName="Helvetica-Bold",
            fontSize=9.5, textColor=COR_TEXTO, leftIndent=22)))
    return KeepTogether(itens)


def montar_progress_bar(spec):
    largura = spec.get("width", 400)
    altura = spec.get("height", 20)
    valor = max(0, min(100, spec.get("value", 0)))
    cor_fundo = colors.HexColor(spec.get("bg", "#E2E8F0"))
    cor_barra = colors.HexColor(spec.get("color", "#4F46E5"))
    drawing = Drawing(largura, altura)
    drawing.add(Rect(0, 0, largura, altura, fillColor=cor_fundo, strokeColor=None, rx=altura / 2, ry=altura / 2))
    largura_preenchida = max(altura, largura * valor / 100)
    drawing.add(Rect(0, 0, largura_preenchida, altura, fillColor=cor_barra, strokeColor=None, rx=altura / 2, ry=altura / 2))
    if spec.get("show_label", True):
        drawing.add(String(largura + 8, altura / 2 - 4, f"{valor}%", fontName="Helvetica-Bold",
                            fontSize=10, fillColor=COR_PRIMARIA))
    return drawing


def montar_badge_lista(items_spec):
    """Lista de badges/tags coloridos lado a lado, tipo etiquetas."""
    itens_flow = []
    for grupo in items_spec:
        cor_bg = colors.HexColor(grupo.get("bg", "#EEF2FF"))
        cor_texto = colors.HexColor(grupo.get("text_color", "#4F46E5"))
        p = Paragraph(grupo.get("text", ""), ParagraphStyle(
            "Badge", fontName="Helvetica-Bold", fontSize=9, textColor=cor_texto, alignment=1))
        tabela = Table([[p]], colWidths=[None])
        tabela.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), cor_bg),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ]))
        itens_flow.append(tabela)
    return itens_flow


def montar_capa(capa_spec, s):
    """Página de capa full-bleed. NÃO usa Drawing como flowable —
    um Drawing do tamanho da página inteira rebenta qualquer Frame
    (que é sempre menor que a página por causa das margens). Em vez
    disso, devolve só os itens de TEXTO da capa (que a Frame absorve
    normalmente) e um dict de instrução para o onPage desenhar o
    fundo colorido diretamente no canvas, ANTES do texto ser
    posicionado por cima."""
    itens = [Spacer(1, 220)]
    p_titulo = Paragraph(capa_spec.get("title", ""), s["capa_titulo"])
    p_titulo._is_cover_flow = True
    itens.append(p_titulo)
    if capa_spec.get("subtitle"):
        p_sub = Paragraph(capa_spec["subtitle"], s["capa_subtitulo"])
        p_sub._is_cover_flow = True
        itens.append(p_sub)
    if capa_spec.get("footer"):
        itens.append(Spacer(1, 340))
        p_rodape = Paragraph(capa_spec["footer"], s["capa_rodape"])
        p_rodape._is_cover_flow = True
        itens.append(p_rodape)
    itens.append(PageBreak())
    return itens


def desenhar_fundo_capa(canvas_obj, doc, capa_spec):
    canvas_obj.saveState()
    largura, altura = doc.pagesize
    cor_fundo = colors.HexColor(capa_spec.get("bg", "#0F172A"))
    canvas_obj.setFillColor(cor_fundo)
    canvas_obj.rect(0, 0, largura, altura, fill=1, stroke=0)
    if capa_spec.get("accent_color"):
        canvas_obj.setFillColor(colors.HexColor(capa_spec["accent_color"]))
        canvas_obj.rect(0, altura / 2 - 2, largura, 4, fill=1, stroke=0)
    canvas_obj.restoreState()


def montar_story_secao(sec, s):
    itens = []

    if sec.get("heading"):
        nivel = sec.get("heading_level", 1)
        estilo_heading = {1: s["heading"], 2: s["heading2"], 3: s["heading3"]}.get(nivel, s["heading"])
        p = Paragraph(sec["heading"], estilo_heading)
        p._toc_heading = (sec["heading"], nivel)
        itens.append(p)

    if sec.get("divider_after_heading"):
        itens.append(linha_divisoria(sec.get("divider_color")))

    if sec.get("runs"):
        itens.append(montar_runs_paragrafo(sec["runs"], s["corpo"]))

    for p in sec.get("paragraphs", []):
        itens.append(Paragraph(p, s["corpo"]))

    for item in sec.get("bullet_list", []):
        itens.append(Paragraph(f"• {item}", s["bullet"]))

    for i, item in enumerate(sec.get("numbered_list", []) or []):
        itens.append(Paragraph(f"{i + 1}. {item}", s["bullet"]))

    if sec.get("quote"):
        itens.append(Spacer(1, 4))
        itens.append(montar_citacao(sec["quote"], s["citacao"]))
        itens.append(Spacer(1, 6))

    if sec.get("callout"):
        itens.append(Spacer(1, 6))
        itens.append(montar_callout(sec["callout"], s["callout"]))
        itens.append(Spacer(1, 8))

    if sec.get("progress_bar"):
        itens.append(Spacer(1, 6))
        itens.append(montar_progress_bar(sec["progress_bar"]))
        itens.append(Spacer(1, 8))

    if sec.get("badges"):
        itens.append(Spacer(1, 4))
        linha_badges = montar_badge_lista(sec["badges"])
        tabela_badges = Table([linha_badges])
        tabela_badges.setStyle(TableStyle([
            ("LEFTPADDING", (0, 0), (-1, -1), 0),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 0),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
        ]))
        itens.append(tabela_badges)
        itens.append(Spacer(1, 8))

    if sec.get("chart"):
        itens.append(Spacer(1, 6))
        itens.append(montar_grafico(sec["chart"]))
        itens.append(Spacer(1, 6))

    if sec.get("charts"):
        for c in sec["charts"]:
            itens.append(Spacer(1, 6))
            itens.append(montar_grafico(c))
            itens.append(Spacer(1, 6))

    if sec.get("shapes"):
        itens.append(Spacer(1, 6))
        itens.append(montar_formas_livres(sec["shapes"]))
        itens.append(Spacer(1, 6))

    if sec.get("qrcode"):
        itens.append(Spacer(1, 6))
        itens.append(montar_qrcode(sec["qrcode"]))
        itens.append(Spacer(1, 6))

    if sec.get("barcode"):
        itens.append(Spacer(1, 6))
        itens.append(montar_barcode128(sec["barcode"]))
        itens.append(Spacer(1, 6))

    if sec.get("image_url"):
        img = baixar_imagem(sec["image_url"], max_width_cm=sec.get("image_width_cm", 15))
        if img:
            itens.append(Spacer(1, 6))
            itens.append(img)
            itens.append(Spacer(1, 6))

    if sec.get("images_grid"):
        linha_atual = []
        linhas = []
        for url in sec["images_grid"]:
            img = baixar_imagem(url, max_width_cm=6)
            if img:
                linha_atual.append(img)
            if len(linha_atual) == 2:
                linhas.append(linha_atual)
                linha_atual = []
        if linha_atual:
            linhas.append(linha_atual)
        for linha in linhas:
            tabela_grid = Table([linha])
            tabela_grid.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ]))
            itens.append(tabela_grid)
            itens.append(Spacer(1, 6))

    if sec.get("table"):
        itens.append(Spacer(1, 6))
        itens.append(montar_tabela(sec["table"]))
        itens.append(Spacer(1, 10))

    if sec.get("keep_together"):
        itens = [KeepTogether(itens)]

    if sec.get("page_break_after"):
        itens.append(PageBreak())

    if sec.get("frame_break_after"):
        itens.append(FrameBreak())

    return itens


class _DocComTOC(BaseDocTemplate):
    def afterFlowable(self, flowable):
        if hasattr(flowable, "_toc_heading"):
            texto, nivel = flowable._toc_heading
            self.notify("TOCEntry", (nivel - 1, texto, self.page))


# ══════════════════════════════════════════════════════════════════
# PDF — função única, absorve tudo do antigo create_pdf_structured
# e aceita também os campos simples do antigo create_pdf (title,
# paragraphs, bullet_list na raiz) por compatibilidade — se vierem
# na raiz sem "sections", são tratados como uma única secção.
# ══════════════════════════════════════════════════════════════════

def create_pdf(params):
    s = get_styles()

    tamanho_nome = params.get("page_size", "A4")
    pagesize = TAMANHOS_PAGINA.get(tamanho_nome, A4)

    num_colunas = params.get("columns", 1)
    usar_toc = bool(params.get("toc"))
    watermark = params.get("watermark")
    footer_spec = params.get("footer")
    faixa_lateral_hex = params.get("side_bar_color")
    header_cor_hex = params.get("header_color")
    capa_spec = params.get("cover")

    sections = params.get("sections")
    if not sections:
        secao_raiz = {
            "paragraphs": params.get("paragraphs", []),
            "bullet_list": params.get("bullet_list", []),
            "numbered_list": params.get("numbered_list", []),
            "runs": params.get("runs"),
            "image_url": params.get("image_url"),
            "images_grid": params.get("images_grid"),
            "table": params.get("table"),
            "chart": params.get("chart"),
            "charts": params.get("charts"),
            "shapes": params.get("shapes"),
            "qrcode": params.get("qrcode"),
            "barcode": params.get("barcode"),
            "callout": params.get("callout"),
            "quote": params.get("quote"),
            "progress_bar": params.get("progress_bar"),
            "badges": params.get("badges"),
            "header_text": params.get("header_text"),
        }
        sections = [secao_raiz]

    header_text = None
    for sec in sections:
        if sec.get("header_text"):
            header_text = sec["header_text"]
            break

    on_page = rodape_e_extras_pagina(header_text, watermark, footer_spec, faixa_lateral_hex, header_cor_hex, capa_spec)

    buffer = io.BytesIO()
    top_margin = 2.6 * cm if header_text else 2.2 * cm

    precisa_doctemplate = num_colunas > 1 or usar_toc or capa_spec

    if precisa_doctemplate:
        doc = _DocComTOC(buffer, pagesize=pagesize,
                          topMargin=top_margin, bottomMargin=2.2 * cm,
                          leftMargin=2 * cm, rightMargin=2 * cm,
                          title=params.get("title", "Documento"))

        largura_disponivel = pagesize[0] - 4 * cm
        if num_colunas > 1:
            largura_coluna = (largura_disponivel - (num_colunas - 1) * 0.6 * cm) / num_colunas
            frames = []
            for i in range(num_colunas):
                x = 2 * cm + i * (largura_coluna + 0.6 * cm)
                frames.append(Frame(x, 2.2 * cm, largura_coluna, pagesize[1] - top_margin - 2.2 * cm,
                                     id=f"col{i}", leftPadding=0, rightPadding=0))
        else:
            frames = [Frame(2 * cm, 2.2 * cm, largura_disponivel, pagesize[1] - top_margin - 2.2 * cm,
                             id="normal", leftPadding=0, rightPadding=0)]

        doc.addPageTemplates([PageTemplate(id="Principal", frames=frames, onPage=on_page)])

        story = []

        if capa_spec:
            story.extend(montar_capa(capa_spec, s))

        story.append(Paragraph(params["title"], s["titulo"]))
        if params.get("subtitle"):
            story.append(Paragraph(params["subtitle"], s["subtitulo"]))
        story.append(linha_divisoria(params.get("title_divider_color")))

        if usar_toc:
            story.append(Paragraph("Índice", s["toc_titulo"]))
            toc = TableOfContents()
            toc.levelStyles = get_toc_styles()
            story.append(toc)
            story.append(PageBreak())

        for sec in sections:
            story.extend(montar_story_secao(sec, s))

        doc.multiBuild(story)

    else:
        doc = SimpleDocTemplate(buffer, pagesize=pagesize,
                                 topMargin=top_margin, bottomMargin=2.2 * cm,
                                 leftMargin=2 * cm, rightMargin=2 * cm,
                                 title=params.get("title", "Documento"))

        story = [Paragraph(params["title"], s["titulo"])]
        if params.get("subtitle"):
            story.append(Paragraph(params["subtitle"], s["subtitulo"]))
        story.append(linha_divisoria(params.get("title_divider_color")))

        for sec in sections:
            story.extend(montar_story_secao(sec, s))

        doc.build(story, onFirstPage=on_page, onLaterPages=on_page)

    pdf_bytes = buffer.getvalue()
    buffer.close()
    return {"pdf_base64": base64.b64encode(pdf_bytes).decode("utf-8")}


# ══════════════════════════════════════════════════════════════════
# DOCX
# ══════════════════════════════════════════════════════════════════

def create_docx(params):
    document = Document()
    document.add_heading(params["title"], level=0)

    if params.get("subtitle"):
        sub = document.add_paragraph(params["subtitle"])
        sub.runs[0].font.size = Pt(13)
        sub.runs[0].font.color.rgb = RGBColor(0x4F, 0x46, 0xE5)

    for sec in params.get("sections", []):
        if sec.get("heading"):
            nivel = sec.get("heading_level", 1)
            document.add_heading(sec["heading"], level=nivel)
        for p in sec.get("paragraphs", []):
            document.add_paragraph(p)
        for item in sec.get("bullet_list", []):
            document.add_paragraph(item, style="List Bullet")
        for item in sec.get("numbered_list", []) or []:
            document.add_paragraph(item, style="List Number")
        if sec.get("image_url"):
            try:
                resp = requests.get(sec["image_url"], timeout=8)
                resp.raise_for_status()
                document.add_picture(io.BytesIO(resp.content), width=Cm(14))
            except Exception:
                pass
        if sec.get("table"):
            t = sec["table"]
            tabela = document.add_table(rows=1, cols=len(t["headers"]))
            tabela.style = "Light Grid Accent 1"
            hdr_cells = tabela.rows[0].cells
            for i, h in enumerate(t["headers"]):
                hdr_cells[i].text = h
            for row in t["rows"]:
                cells = tabela.add_row().cells
                for i, val in enumerate(row):
                    cells[i].text = str(val)

    buffer = io.BytesIO()
    document.save(buffer)
    docx_bytes = buffer.getvalue()
    buffer.close()
    return {"docx_base64": base64.b64encode(docx_bytes).decode("utf-8")}


# ══════════════════════════════════════════════════════════════════
# XLSX — openpyxl, múltiplas sheets, gráficos nativos do Excel,
# formatação condicional, tabelas estilizadas, congelar painéis,
# larguras/alturas customizadas, imagens embutidas.
# ══════════════════════════════════════════════════════════════════

def _resolver_cor_xlsx(hex_str):
    return hex_str.lstrip("#").upper() if hex_str else None


def _aplicar_estilo_cabecalho_xlsx(ws, linha_idx, num_cols, bg_hex="0F172A", texto_hex="FFFFFF"):
    fill = PatternFill(start_color=_resolver_cor_xlsx(bg_hex), end_color=_resolver_cor_xlsx(bg_hex), fill_type="solid")
    fonte = Font(bold=True, color=_resolver_cor_xlsx(texto_hex), size=11)
    for col_idx in range(1, num_cols + 1):
        celula = ws.cell(row=linha_idx, column=col_idx)
        celula.fill = fill
        celula.font = fonte
        celula.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)


def _aplicar_bordas_xlsx(ws, linha_inicio, linha_fim, col_inicio, col_fim, cor_hex="DDDDDD"):
    borda = Border(
        left=Side(style="thin", color=_resolver_cor_xlsx(cor_hex)),
        right=Side(style="thin", color=_resolver_cor_xlsx(cor_hex)),
        top=Side(style="thin", color=_resolver_cor_xlsx(cor_hex)),
        bottom=Side(style="thin", color=_resolver_cor_xlsx(cor_hex)),
    )
    for linha in range(linha_inicio, linha_fim + 1):
        for col in range(col_inicio, col_fim + 1):
            ws.cell(row=linha, column=col).border = borda


def _montar_grafico_xlsx(ws, chart_spec, linha_dados_inicio, linha_dados_fim, col_labels, col_dados_inicio, col_dados_fim, ancora):
    tipo = chart_spec.get("type", "bar")

    if tipo == "pie":
        chart = PieChart()
    elif tipo == "line":
        chart = LineChart()
    elif tipo == "scatter":
        chart = ScatterChart()
    else:
        chart = BarChart()
        chart.type = "col"
        if chart_spec.get("horizontal"):
            chart.type = "bar"

    chart.title = chart_spec.get("title", "")
    chart.style = chart_spec.get("style_id", 10)
    chart.width = chart_spec.get("width_cm", 16)
    chart.height = chart_spec.get("height_cm", 9)

    if tipo != "scatter":
        dados = Reference(ws, min_col=col_dados_inicio, max_col=col_dados_fim,
                           min_row=linha_dados_inicio - 1, max_row=linha_dados_fim)
        categorias = Reference(ws, min_col=col_labels, max_col=col_labels,
                                min_row=linha_dados_inicio, max_row=linha_dados_fim)
        chart.add_data(dados, titles_from_data=True)
        chart.set_categories(categorias)

    cores_serie = chart_spec.get("colors") or PALETA_SEM_HASH
    if tipo == "pie" and hasattr(chart, "series") and chart.series:
        from openpyxl.chart.marker import DataPoint
        pontos = []
        n_fatias = linha_dados_fim - linha_dados_inicio + 1
        for i in range(n_fatias):
            dp = DataPoint(idx=i)
            dp.graphicalProperties.solidFill = cores_serie[i % len(cores_serie)]
            pontos.append(dp)
        chart.series[0].data_points = pontos
    elif hasattr(chart, "series"):
        for i, serie in enumerate(chart.series):
            cor = cores_serie[i % len(cores_serie)]
            if tipo == "line":
                serie.graphicalProperties.line.solidFill = cor
                serie.graphicalProperties.line.width = 22000
                serie.marker.symbol = "circle"
                serie.marker.size = 6
                serie.marker.graphicalProperties.solidFill = cor
            else:
                serie.graphicalProperties.solidFill = cor

    if chart_spec.get("legend_position"):
        chart.legend.position = chart_spec["legend_position"]
    if chart_spec.get("show_legend") is False:
        chart.legend = None

    ws.add_chart(chart, ancora)
    return chart


def _preencher_sheet_xlsx(wb, sheet_spec, is_primeira=False):
    nome = sheet_spec.get("name", "Sheet1")[:31]
    if is_primeira:
        ws = wb.active
        ws.title = nome
    else:
        ws = wb.create_sheet(title=nome)

    if sheet_spec.get("tab_color"):
        ws.sheet_properties.tabColor = _resolver_cor_xlsx(sheet_spec["tab_color"])

    headers = sheet_spec.get("headers", [])
    rows = sheet_spec.get("rows", [])
    linha_titulo_offset = 0

    if sheet_spec.get("title"):
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(1, len(headers)))
        celula_titulo = ws.cell(row=1, column=1, value=sheet_spec["title"])
        celula_titulo.font = Font(bold=True, size=16, color=_resolver_cor_xlsx(sheet_spec.get("title_color", "0F172A")))
        celula_titulo.alignment = Alignment(horizontal="left", vertical="center")
        ws.row_dimensions[1].height = 28
        linha_titulo_offset = 2

    linha_cabecalho = 1 + linha_titulo_offset
    for col_idx, h in enumerate(headers, start=1):
        ws.cell(row=linha_cabecalho, column=col_idx, value=h)
    if headers:
        _aplicar_estilo_cabecalho_xlsx(
            ws, linha_cabecalho, len(headers),
            bg_hex=sheet_spec.get("header_bg", "0F172A"),
            texto_hex=sheet_spec.get("header_text_color", "FFFFFF"),
        )
        ws.row_dimensions[linha_cabecalho].height = 22

    linha_dados_inicio = linha_cabecalho + 1
    cor_zebra = _resolver_cor_xlsx(sheet_spec.get("zebra_color", "F8FAFC"))
    fill_zebra = PatternFill(start_color=cor_zebra, end_color=cor_zebra, fill_type="solid")

    for i, row in enumerate(rows):
        linha_atual = linha_dados_inicio + i
        for col_idx, val in enumerate(row, start=1):
            celula = ws.cell(row=linha_atual, column=col_idx, value=val)
            if isinstance(val, (int, float)) and sheet_spec.get("number_format"):
                celula.number_format = sheet_spec["number_format"]
            if i % 2 == 1 and sheet_spec.get("zebra", True):
                celula.fill = fill_zebra
            celula.alignment = Alignment(vertical="center")

    linha_dados_fim = linha_dados_inicio + len(rows) - 1

    if headers and rows:
        _aplicar_bordas_xlsx(ws, linha_cabecalho, linha_dados_fim, 1, len(headers))

    if sheet_spec.get("as_excel_table") and headers and rows:
        ref = f"{get_column_letter(1)}{linha_cabecalho}:{get_column_letter(len(headers))}{linha_dados_fim}"
        nome_tabela = "".join(ch for ch in nome if ch.isalnum()) or "Tabela1"
        tabela_excel = XlsxTable(displayName=f"Tbl{nome_tabela}{abs(hash(ref)) % 10000}", ref=ref)
        tabela_excel.tableStyleInfo = TableStyleInfo(
            name=sheet_spec.get("table_style", "TableStyleMedium9"),
            showRowStripes=True, showFirstColumn=False,
        )
        ws.add_table(tabela_excel)

    col_widths = sheet_spec.get("col_widths")
    if col_widths:
        for i, w in enumerate(col_widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w
    elif headers:
        for i, h in enumerate(headers, start=1):
            maior = max([len(str(h))] + [len(str(r[i - 1])) for r in rows if i - 1 < len(r)] or [10])
            ws.column_dimensions[get_column_letter(i)].width = min(40, max(10, maior + 4))

    if sheet_spec.get("freeze_header") and headers:
        ws.freeze_panes = ws.cell(row=linha_dados_inicio, column=1).coordinate

    if sheet_spec.get("conditional_color_scale") and headers and rows:
        col_idx = sheet_spec["conditional_color_scale"].get("col", 1)
        col_letra = get_column_letter(col_idx)
        rng = f"{col_letra}{linha_dados_inicio}:{col_letra}{linha_dados_fim}"
        regra = ColorScaleRule(
            start_type="min", start_color="F8696B",
            mid_type="percentile", mid_value=50, mid_color="FFEB84",
            end_type="max", end_color="63BE7B",
        )
        ws.conditional_formatting.add(rng, regra)

    if sheet_spec.get("data_bars") and headers and rows:
        col_idx = sheet_spec["data_bars"].get("col", 1)
        col_letra = get_column_letter(col_idx)
        rng = f"{col_letra}{linha_dados_inicio}:{col_letra}{linha_dados_fim}"
        regra = DataBarRule(
            start_type="min", end_type="max",
            color=_resolver_cor_xlsx(sheet_spec["data_bars"].get("color", "4F46E5")),
        )
        ws.conditional_formatting.add(rng, regra)

    if sheet_spec.get("chart") and headers and rows:
        col_dados_inicio = sheet_spec["chart"].get("data_col_start", 2)
        col_dados_fim = sheet_spec["chart"].get("data_col_end", len(headers))
        _montar_grafico_xlsx(
            ws, sheet_spec["chart"], linha_dados_inicio, linha_dados_fim,
            col_labels=1, col_dados_inicio=col_dados_inicio, col_dados_fim=col_dados_fim,
            ancora=sheet_spec["chart"].get("anchor", f"{get_column_letter(len(headers) + 2)}{linha_cabecalho}"),
        )

    for c in sheet_spec.get("charts", []) or []:
        col_dados_inicio = c.get("data_col_start", 2)
        col_dados_fim = c.get("data_col_end", len(headers))
        _montar_grafico_xlsx(
            ws, c, linha_dados_inicio, linha_dados_fim,
            col_labels=1, col_dados_inicio=col_dados_inicio, col_dados_fim=col_dados_fim,
            ancora=c.get("anchor", f"{get_column_letter(len(headers) + 2)}{linha_cabecalho}"),
        )

    if sheet_spec.get("image_url"):
        img_bytes = baixar_bytes_imagem(sheet_spec["image_url"])
        if img_bytes:
            xlsx_img = XlsxImage(img_bytes)
            largura_col_img = sheet_spec.get("image_width_px", 300)
            xlsx_img.width = largura_col_img
            xlsx_img.height = int(largura_col_img * 0.6)
            ws.add_image(xlsx_img, sheet_spec.get("image_anchor", "A1"))

    if sheet_spec.get("notes"):
        linha_notas = linha_dados_fim + 2
        for i, nota in enumerate(sheet_spec["notes"]):
            celula = ws.cell(row=linha_notas + i, column=1, value=nota)
            celula.font = Font(italic=True, size=9, color=_resolver_cor_xlsx("475569"))

    return ws


def create_xlsx(params):
    wb = Workbook()

    sheets = params.get("sheets")
    if not sheets:
        sheets = [{
            "name": params.get("sheet_name", "Dados"),
            "title": params.get("title"),
            "headers": params.get("headers", []),
            "rows": params.get("rows", []),
            "chart": params.get("chart"),
            "charts": params.get("charts"),
            "col_widths": params.get("col_widths"),
            "freeze_header": params.get("freeze_header", True),
            "as_excel_table": params.get("as_excel_table"),
            "conditional_color_scale": params.get("conditional_color_scale"),
            "data_bars": params.get("data_bars"),
            "zebra": params.get("zebra", True),
            "notes": params.get("notes"),
            "image_url": params.get("image_url"),
        }]

    for i, sheet_spec in enumerate(sheets):
        _preencher_sheet_xlsx(wb, sheet_spec, is_primeira=(i == 0))

    buffer = io.BytesIO()
    wb.save(buffer)
    xlsx_bytes = buffer.getvalue()
    buffer.close()
    return {"xlsx_base64": base64.b64encode(xlsx_bytes).decode("utf-8")}


# ══════════════════════════════════════════════════════════════════
# PPTX — python-pptx, layouts completos, gráficos nativos, tabelas,
# imagens, formas, notas do orador, temas de cor por slide.
# ══════════════════════════════════════════════════════════════════

def _cor_pptx(hex_str):
    h = hex_str.lstrip("#")
    return PptxRGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _limpar_slide_layout_em_branco(prs):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _adicionar_fundo_solido(slide, hex_cor):
    fundo = slide.background
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = _cor_pptx(hex_cor)


def _add_titulo_slide(slide, texto, cor_hex="0F172A", tamanho=32, top_in=0.4, left_in=0.6, width_in=8.8, bold=True):
    caixa = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(1.0))
    tf = caixa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = PptxPt(tamanho)
    run.font.bold = bold
    run.font.color.rgb = _cor_pptx(cor_hex)
    run.font.name = "Calibri"
    return caixa


def _add_subtitulo_slide(slide, texto, cor_hex="4F46E5", tamanho=16, top_in=1.15, left_in=0.6, width_in=8.8):
    caixa = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.6))
    tf = caixa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = PptxPt(tamanho)
    run.font.color.rgb = _cor_pptx(cor_hex)
    run.font.name = "Calibri"
    return caixa


def _add_bullets_slide(slide, bullets, top_in=1.9, left_in=0.7, width_in=8.6, height_in=5.0,
                        cor_hex="475569", tamanho=16, numbered=False):
    caixa = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = caixa.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefixo = f"{i + 1}. " if numbered else "•  "
        run = p.add_run()
        run.text = prefixo + item
        run.font.size = PptxPt(tamanho)
        run.font.color.rgb = _cor_pptx(cor_hex)
        run.font.name = "Calibri"
        p.space_after = PptxPt(10)
    return caixa


def _add_imagem_slide(slide, url, left_in, top_in, width_in=None, height_in=None):
    img_bytes = baixar_bytes_imagem(url)
    if not img_bytes:
        return None
    kwargs = {}
    if width_in:
        kwargs["width"] = Inches(width_in)
    if height_in:
        kwargs["height"] = Inches(height_in)
    return slide.shapes.add_picture(img_bytes, Inches(left_in), Inches(top_in), **kwargs)


def _add_tabela_slide(slide, table_spec, left_in=0.6, top_in=1.8, width_in=8.8, height_in=4.5):
    headers = table_spec["headers"]
    rows = table_spec["rows"]
    n_linhas = len(rows) + 1
    n_colunas = len(headers)

    shape = slide.shapes.add_table(n_linhas, n_colunas, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tabela = shape.table

    cor_cabecalho = _cor_pptx(table_spec.get("header_bg", "0F172A"))
    for col_idx, h in enumerate(headers):
        celula = tabela.cell(0, col_idx)
        celula.text = str(h)
        celula.fill.solid()
        celula.fill.fore_color.rgb = cor_cabecalho
        for p in celula.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = PptxPt(table_spec.get("header_font_size", 13))
            p.font.color.rgb = _cor_pptx(table_spec.get("header_text_color", "FFFFFF"))

    cor_zebra = _cor_pptx(table_spec.get("zebra_color", "F8FAFC"))
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, val in enumerate(row):
            celula = tabela.cell(row_idx, col_idx)
            celula.text = str(val)
            if row_idx % 2 == 0 and table_spec.get("zebra", True):
                celula.fill.solid()
                celula.fill.fore_color.rgb = cor_zebra
            for p in celula.text_frame.paragraphs:
                p.font.size = PptxPt(table_spec.get("font_size", 11))
                p.font.color.rgb = _cor_pptx(table_spec.get("text_color", "475569"))

    return shape


def _add_grafico_slide(slide, chart_spec, left_in=0.8, top_in=1.9, width_in=8.4, height_in=4.6):
    tipo_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_tipo = tipo_map.get(chart_spec.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    tipo = chart_spec.get("type", "bar")

    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.get("labels", [])
    for serie in chart_spec.get("series", []):
        chart_data.add_series(serie.get("name", "Série"), serie.get("data", []))

    grafico_shape = slide.shapes.add_chart(xl_tipo, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in), chart_data)
    chart = grafico_shape.chart

    if chart_spec.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec["title"]
    else:
        chart.has_title = False

    chart.has_legend = len(chart_spec.get("series", [])) > 1 and chart_spec.get("show_legend", True)
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    cores = chart_spec.get("colors") or PALETA
    try:
        plot = chart.plots[0]
        if tipo in ("pie", "doughnut"):
            for i, ponto in enumerate(plot.series[0].points):
                ponto.format.fill.solid()
                ponto.format.fill.fore_color.rgb = _cor_pptx(cores[i % len(cores)])
        else:
            for i, ponto_serie in enumerate(plot.series):
                cor = cores[i % len(cores)]
                ponto_serie.format.fill.solid()
                ponto_serie.format.fill.fore_color.rgb = _cor_pptx(cor)
    except Exception:
        pass

    return grafico_shape


def _add_forma_slide(slide, shape_spec):
    tipo_map = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOCELES_TRIANGLE,
        "arrow_right": MSO_SHAPE.RIGHT_ARROW,
        "star": MSO_SHAPE.STAR_5_POINT,
        "chevron": MSO_SHAPE.CHEVRON,
        "hexagon": MSO_SHAPE.HEXAGON,
    }
    mso_tipo = tipo_map.get(shape_spec.get("type", "rect"), MSO_SHAPE.RECTANGLE)
    forma = slide.shapes.add_shape(
        mso_tipo,
        Inches(shape_spec.get("left", 1)), Inches(shape_spec.get("top", 1)),
        Inches(shape_spec.get("width", 2)), Inches(shape_spec.get("height", 1)),
    )
    if shape_spec.get("fill"):
        forma.fill.solid()
        forma.fill.fore_color.rgb = _cor_pptx(shape_spec["fill"])
    else:
        forma.fill.background()
    if shape_spec.get("stroke"):
        forma.line.color.rgb = _cor_pptx(shape_spec["stroke"])
        forma.line.width = PptxPt(shape_spec.get("stroke_width", 1))
    else:
        forma.line.fill.background()
    if shape_spec.get("text"):
        forma.text_frame.text = shape_spec["text"]
        for p in forma.text_frame.paragraphs:
            p.font.size = PptxPt(shape_spec.get("font_size", 14))
            p.font.bold = shape_spec.get("bold", False)
            p.font.color.rgb = _cor_pptx(shape_spec.get("text_color", "FFFFFF"))
            p.alignment = PP_ALIGN.CENTER
        forma.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return forma


def _montar_slide(prs, layout_branco, slide_spec, defaults):
    slide = prs.slides.add_slide(layout_branco)

    bg = slide_spec.get("bg", defaults.get("bg"))
    if bg:
        _adicionar_fundo_solido(slide, bg)

    if slide_spec.get("type") == "cover":
        _add_titulo_slide(slide, slide_spec.get("title", ""),
                           cor_hex=slide_spec.get("title_color", "FFFFFF"),
                           tamanho=slide_spec.get("title_size", 40),
                           top_in=2.6, left_in=0.6, width_in=8.8)
        if slide_spec.get("subtitle"):
            _add_subtitulo_slide(slide, slide_spec["subtitle"],
                                  cor_hex=slide_spec.get("subtitle_color", "E2E8F0"),
                                  tamanho=18, top_in=3.6)
        if slide_spec.get("author") or slide_spec.get("date"):
            rodape = " | ".join(filter(None, [slide_spec.get("author"), slide_spec.get("date")]))
            _add_subtitulo_slide(slide, rodape, cor_hex=slide_spec.get("footer_color", "94A3B8"), tamanho=12, top_in=6.7)
        return slide

    if slide_spec.get("type") == "section":
        _add_titulo_slide(slide, slide_spec.get("title", ""),
                           cor_hex=slide_spec.get("title_color", "FFFFFF"),
                           tamanho=slide_spec.get("title_size", 34),
                           top_in=3.1, left_in=0.6, width_in=8.8)
        if slide_spec.get("subtitle"):
            _add_subtitulo_slide(slide, slide_spec["subtitle"],
                                  cor_hex=slide_spec.get("subtitle_color", "E2E8F0"),
                                  tamanho=15, top_in=4.0)
        return slide

    if slide_spec.get("heading"):
        _add_titulo_slide(slide, slide_spec["heading"],
                           cor_hex=slide_spec.get("title_color", defaults.get("title_color", "0F172A")),
                           tamanho=slide_spec.get("title_size", 28))

    if slide_spec.get("subheading"):
        _add_subtitulo_slide(slide, slide_spec["subheading"],
                              cor_hex=slide_spec.get("subtitle_color", defaults.get("subtitle_color", "4F46E5")))

    if slide_spec.get("bullets"):
        _add_bullets_slide(slide, slide_spec["bullets"],
                            cor_hex=slide_spec.get("text_color", defaults.get("text_color", "475569")),
                            numbered=slide_spec.get("numbered", False))

    if slide_spec.get("image_url") and not slide_spec.get("image_left"):
        _add_imagem_slide(slide, slide_spec["image_url"], left_in=5.3, top_in=1.9, width_in=4.0)

    if slide_spec.get("image_left"):
        _add_imagem_slide(slide, slide_spec["image_left"], left_in=0.6, top_in=1.9, width_in=4.0)

    if slide_spec.get("images_grid"):
        posicoes = [(0.6, 1.9), (5.3, 1.9), (0.6, 4.6), (5.3, 4.6)]
        for url, (lx, ty) in zip(slide_spec["images_grid"], posicoes):
            _add_imagem_slide(slide, url, left_in=lx, top_in=ty, width_in=3.9)

    if slide_spec.get("table"):
        _add_tabela_slide(slide, slide_spec["table"])

    if slide_spec.get("chart"):
        _add_grafico_slide(slide, slide_spec["chart"])

    for forma_spec in slide_spec.get("shapes", []) or []:
        _add_forma_slide(slide, forma_spec)

    if slide_spec.get("quote"):
        caixa = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(8.0), Inches(2.0))
        tf = caixa.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f'"{slide_spec["quote"].get("text", "")}"'
        run.font.italic = True
        run.font.size = PptxPt(24)
        run.font.color.rgb = _cor_pptx(slide_spec.get("title_color", "0F172A"))
        if slide_spec["quote"].get("author"):
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = f'— {slide_spec["quote"]["author"]}'
            run2.font.bold = True
            run2.font.size = PptxPt(15)
            run2.font.color.rgb = _cor_pptx("4F46E5")

    if slide_spec.get("footer_text"):
        caixa = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(6), Inches(0.4))
        p = caixa.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = slide_spec["footer_text"]
        run.font.size = PptxPt(9)
        run.font.color.rgb = _cor_pptx("94A3B8")

    if slide_spec.get("slide_number", True):
        caixa = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(0.6), Inches(0.4))
        p = caixa.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(slide_spec.get("_numero", ""))
        run.font.size = PptxPt(9)
        run.font.color.rgb = _cor_pptx("94A3B8")
        p.alignment = PP_ALIGN.RIGHT

    if slide_spec.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_spec["speaker_notes"]

    return slide


def create_pptx(params):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    layout_branco = _limpar_slide_layout_em_branco(prs)

    defaults = params.get("theme", {})

    slides = params.get("slides")
    if not slides:
        slides = [{
            "heading": params.get("title", "Slide"),
            "bullets": params.get("bullets", []),
        }]

    if params.get("title") and params.get("cover", True) and slides and slides[0].get("type") != "cover":
        slides = [{
            "type": "cover",
            "title": params["title"],
            "subtitle": params.get("subtitle"),
            "author": params.get("author"),
            "date": params.get("date"),
            "bg": defaults.get("cover_bg", "0F172A"),
        }] + slides

    for i, slide_spec in enumerate(slides, start=1):
        slide_spec["_numero"] = i
        _montar_slide(prs, layout_branco, slide_spec, defaults)

    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    buffer.close()
    return {"pptx_base64": base64.b64encode(pptx_bytes).decode("utf-8")}


# ══════════════════════════════════════════════════════════════════
# Mindmap PNG (mantido do original)
# ══════════════════════════════════════════════════════════════════

def _draw_mindmap_node(drawing, node, x, y, depth, color_idx, max_width):
    label = node.get("label", "")
    bg_color = colors.HexColor(PALETA[color_idx % len(PALETA)]) if depth > 0 else COR_PRIMARIA
    box_width = min(max_width - x, 20 + len(label) * 6.5)
    box_height = 22

    rect = Rect(x, y - box_height, box_width, box_height, fillColor=bg_color, strokeColor=None, rx=6, ry=6)
    drawing.add(rect)
    text = String(x + 10, y - box_height + 7, label, fontSize=max(8, 12 - depth), fillColor=colors.white, fontName="Helvetica-Bold")
    drawing.add(text)

    current_y = y - box_height - 10
    children = node.get("children", [])
    for i, child in enumerate(children):
        current_y = _draw_mindmap_node(drawing, child, x + 24, current_y, depth + 1, color_idx + i + 1, max_width)
    return current_y


def generate_mindmap(params):
    root = params["root"]
    width, height = 960, 720
    drawing = Drawing(width, height)
    drawing.add(Rect(0, 0, width, height, fillColor=colors.white, strokeColor=None))
    _draw_mindmap_node(drawing, root, 20, height - 20, 0, 0, width)

    png_buffer = io.BytesIO()
    renderPM.drawToFile(drawing, png_buffer, fmt="PNG")
    png_bytes = png_buffer.getvalue()
    png_buffer.close()
    return {"png_base64": base64.b64encode(png_bytes).decode("utf-8")}


if __name__ == "__main__":
    func_name = sys.argv[1]
    params = json.loads(sys.argv[2])

    funcoes = {
        "create_pdf": create_pdf,
        "create_docx": create_docx,
        "create_xlsx": create_xlsx,
        "create_pptx": create_pptx,
        "generate_mindmap": generate_mindmap,
    }

    if func_name not in funcoes:
        print(json.dumps({"error": f"Função '{func_name}' não encontrada"}))
        sys.exit(1)

    try:
        resultado = funcoes[func_name](params)
        print(json.dumps(resultado))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)